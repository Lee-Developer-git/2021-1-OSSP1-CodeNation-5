# -*- coding: utf-8 -*- 
from tika import parser
from pptx import Presentation
import olefile
import sys
import json
import docx2txt
from krwordrank.hangle import normalize;
from collections import defaultdict
import math
import numpy as np

# https://lovit.github.io/nlp/2018/04/16/krwordrank/ 
def hits(graph, beta, max_iter=50, bias=None, verbose=True, 
    sum_weight=100, number_of_nodes=None, converge=0.001):

    if not bias:
        bias = {}
    if not number_of_nodes:
        number_of_nodes = max(len(graph), len(bias))

    if number_of_nodes <= 1:
        raise ValueError(
            'The graph should consist of at least two nodes\n',
            'The node size of inserted graph is %d' % number_of_nodes
        )

    dw = sum_weight / number_of_nodes
    rank = {node:dw for node in graph.keys()}

    for num_iter in range(1, max_iter + 1):
        rank_ = _update(rank, graph, bias, dw, beta)
        diff = sum((abs(w - rank.get(n, 0)) for n, w in rank_.items()))
        rank = rank_

        if diff < sum_weight * converge:
            if verbose:
                pass
            break

        if verbose:
            pass

    if verbose:
        pass

    return rank

def _update(rank, graph, bias, dw, beta):
    rank_new = {}
    for to_node, from_dict in graph.items():
        rank_new[to_node] = sum([w * rank[from_node] for from_node, w in from_dict.items()])
        rank_new[to_node] = beta * rank_new[to_node] + (1 - beta) * bias.get(to_node, dw)
    return rank_new


def summarize_with_keywords(texts, num_keywords=100, stopwords=None, min_count=5,
    max_length=10, beta=0.85, max_iter=10, num_rset=-1, verbose=False):

    # train KR-WordRank
    wordrank_extractor = KRWordRank(
        min_count = min_count,
        max_length = max_length,
        verbose = verbose
        )

    keywords, rank, graph = wordrank_extractor.extract(texts,
        beta, max_iter, num_rset=num_rset)

    # stopword filtering
    if stopwords is None:
        stopwords = {}
    keywords = {word:r for word, r in keywords.items() if not (word in stopwords)}

    # top rank filtering
    if num_keywords > 0:
        keywords = {word:r for word, r in sorted(keywords.items(), key=lambda x:-x[1])[:num_keywords]}

    return keywords


class KRWordRank:
    def __init__(self, min_count=5, max_length=10, verbose=False):
        self.min_count = min_count
        self.max_length = max_length
        self.verbose = verbose
        self.sum_weight = 1
        self.vocabulary = {}
        self.index2vocab = []

    def scan_vocabs(self, docs):
        self.vocabulary = {}
        if self.verbose:
            pass

        counter = {}
        for doc in docs:

            for token in doc.split():
                len_token = len(token)
                counter[(token, 'L')] = counter.get((token, 'L'), 0) + 1

                for e in range(1, min(len(token), self.max_length)):
                    if (len_token - e) > self.max_length:
                        continue

                    l_sub = (token[:e], 'L')
                    r_sub = (token[e:], 'R')
                    counter[l_sub] = counter.get(l_sub, 0) + 1
                    counter[r_sub] = counter.get(r_sub, 0) + 1

        counter = {token:freq for token, freq in counter.items() if freq >= self.min_count}
        for token, _ in sorted(counter.items(), key=lambda x:x[1], reverse=True):
            self.vocabulary[token] = len(self.vocabulary)

        self._build_index2vocab()

        if self.verbose:
            pass
        return counter

    def _build_index2vocab(self):
        self.index2vocab = [vocab for vocab, index in sorted(self.vocabulary.items(), key=lambda x:x[1])]
        self.sum_weight = len(self.index2vocab)
    
    def extract(self, docs, beta=0.85, max_iter=10, num_keywords=-1,
        num_rset=-1, vocabulary=None, bias=None, rset=None):
        rank, graph = self.train(docs, beta, max_iter, vocabulary, bias)

        lset = {self.int2token(idx)[0]:r for idx, r in rank.items() if self.int2token(idx)[1] == 'L'}
        if not rset:
            rset = {self.int2token(idx)[0]:r for idx, r in rank.items() if self.int2token(idx)[1] == 'R'}

        if num_rset > 0:
            rset = {token:r for token, r in sorted(rset.items(), key=lambda x:-x[1])[:num_rset]}

        keywords = self._select_keywords(lset, rset)
        keywords = self._filter_compounds(keywords)
        keywords = self._filter_subtokens(keywords)

        if num_keywords > 0:
            keywords = {token:r for token, r in sorted(keywords.items(), key=lambda x:-x[1])[:num_keywords]}

        return keywords, rank, graph

    def _select_keywords(self, lset, rset):
        keywords = {}
        for word, r in sorted(lset.items(), key=lambda x:x[1], reverse=True):
            len_word = len(word)
            if len_word == 1:
                continue

            is_compound = False
            for e in range(2, len_word):
                if (word[:e] in keywords) and (word[:e] in rset):
                    is_compound = True
                    break

            if not is_compound:
                keywords[word] = r

        return keywords

    def _filter_compounds(self, keywords):
        keywords_= {}
        for word, r in sorted(keywords.items(), key=lambda x:x[1], reverse=True):
            len_word = len(word)

            if len_word <= 2:
                keywords_[word] = r
                continue

            if len_word == 3:
                if word[:2] in keywords_:
                    continue

            is_compound = False
            for e in range(2, len_word - 1):
                # fixed. comment from Y. cho
                if (word[:e] in keywords) and (word[e:] in keywords):
                    is_compound = True
                    break

            if not is_compound:
                keywords_[word] = r

        return keywords_

    def _filter_subtokens(self, keywords):
        subtokens = set()
        keywords_ = {}

        for word, r in sorted(keywords.items(), key=lambda x:x[1], reverse=True):
            subs = {word[:e] for e in range(2, len(word)+1)}

            is_subtoken = False
            for sub in subs:
                if sub in subtokens:
                    is_subtoken = True
                    break

            if not is_subtoken:
                keywords_[word] = r
                subtokens.update(subs)

        return keywords_

    def train(self, docs, beta=0.85, max_iter=10, vocabulary=None, bias=None):
        if (not vocabulary) and (not self.vocabulary):
            self.scan_vocabs(docs)
        elif (not vocabulary):
            self.vocabulary = vocabulary
            self._build_index2vocab()

        graph = self._construct_word_graph(docs)

        rank = hits(graph, beta, max_iter, bias,
                    sum_weight=self.sum_weight,
                    number_of_nodes=len(self.vocabulary),
                    verbose=self.verbose
                    )

        return rank, graph

    def token2int(self, token):
        return self.vocabulary.get(token, -1)

    def int2token(self, index):
        return self.index2vocab[index] if (0 <= index < len(self.index2vocab)) else None

    def _construct_word_graph(self, docs):
        def normalize(graph):
            graph_ = defaultdict(lambda: defaultdict(lambda: 0))
            for from_, to_dict in graph.items():
                sum_ = sum(to_dict.values())
                for to_, w in to_dict.items():
                    graph_[to_][from_] = w / sum_
            graph_ = {t:dict(fd) for t, fd in graph_.items()}
            return graph_

        graph = defaultdict(lambda: defaultdict(lambda: 0))
        for doc in docs:

            tokens = doc.split()

            if not tokens:
                continue

            links = []
            for token in tokens:
                links += self._intra_link(token)

            if len(tokens) > 1:
                tokens = [tokens[-1]] + tokens + [tokens[0]]
                links += self._inter_link(tokens)

            links = self._check_token(links)
            if not links:
                continue

            links = self._encode_token(links)
            for l_node, r_node in links:
                graph[l_node][r_node] += 1
                graph[r_node][l_node] += 1

        # reverse for inbound graph. but it normalized with sum of outbound weight
        graph = normalize(graph)
        return graph

    def _intra_link(self, token):
        links = []
        len_token = len(token)
        for e in range(1, min(len_token, 10)):
            if (len_token - e) > self.max_length:
                continue
            links.append( ((token[:e], 'L'), (token[e:], 'R')) )
        return links

    def _inter_link(self, tokens):
        def rsub_to_token(t_left, t_curr):
            return [((t_left[-b:], 'R'), (t_curr, 'L')) for b in range(1, min(10, len(t_left)))]
        def token_to_lsub(t_curr, t_rigt):
            return [((t_curr, 'L'), (t_rigt[:e], 'L')) for e in range(1, min(10, len(t_rigt)))]

        links = []
        for i in range(1, len(tokens)-1):
            links += rsub_to_token(tokens[i-1], tokens[i])
            links += token_to_lsub(tokens[i], tokens[i+1])
        return links

    def _check_token(self, token_list):
        return [(token[0], token[1]) for token in token_list if (token[0] in self.vocabulary and token[1] in self.vocabulary)]

    def _encode_token(self, token_list):
        return [(self.vocabulary[token[0]],self.vocabulary[token[1]]) for token in token_list]

def get_path(path):
    return path

# pdf to txt
def pdf_to_txt(path):
    pdf_path = path
    raw_pdf = parser.from_file(pdf_path) 
    result = raw_pdf['content'] 
    result = result.strip()
    #print(result)
    return result


# ppt to txt 
def ppt_to_txt(path):
    ppt_path = Presentation(path)
    list = []
    for slide in ppt_path.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                list.append(paragraph.text)
    result = " ".join(list)
    #print(result)
    return result


# hwp to txt ---> Prvtxt만 가능
def hwp_to_txt(path):
    hwp_file = olefile.OleFileIO(path)
    hwp_txt = hwp_file.openstream('Prvtext').read()
    result = hwp_txt.decode('UTF=16')
    #print (result)
    return result


# word to txt
def word_to_txt(path):
    result = docx2txt.process(path)
    #print(result)
    return result


# keyword extraction from txt
def keyword_extraction(txt):
    list_str = txt.split()
    list_file = []
    for line in list_str:
        list_file.append(line)

    texts = list_file
    texts = [normalize(text, english=True, number=True) for text in texts]

    wordrank_extractor = KRWordRank(
        min_count = 5, # 단어의 최소 출현 빈도수 (그래프 생성 시)
        max_length = 10, # 단어의 최대 길이
        verbose = True
        )

    beta = 0.85    # PageRank의 decaying factor beta
    max_iter = 10

    keywords, rank, graph = wordrank_extractor.extract(texts, beta, max_iter)

    for word, r in sorted(keywords.items(), key=lambda x:x[1], reverse=True)[:5]: 
        result_list.append(word)
        #print(word)
        #print('%8s:\t%.4f' % (word, r))


if __name__ == '__main__': 
     p = get_path(sys.argv[1])




if ("pdf" in p):
    result = pdf_to_txt(p)

elif ("pptx" in p):
    result = ppt_to_txt(p) 

elif ("hwp" in p):
    result = hwp_to_txt(p)

elif ("docx" in p):
    result = word_to_txt(p)

else:
    result = "error"


result_list = []
if (result != "error"):
    keyword_extraction(result)


output = dict(zip(range(1, len(result_list) + 1), result_list))
json.dumps(output)
print(output)

# file_path = "./"+"1"+".json"


# with open(file_path, 'w', encoding='UTF-8-sig') as outfile:
#     json.dump(output, outfile, ensure_ascii=False)
